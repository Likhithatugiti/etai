"""
Universal Document Ingestor
Supports: PDF, DOCX, PPTX, CSV, XLSX, TXT
Returns a list of chunk dicts ready for KG + vector store.
"""

import os
import uuid
from typing import List, Dict, Any

CHUNK_SIZE = 500   # characters
CHUNK_OVERLAP = 80


def _chunk_text(text: str, doc_id: str, doc_type: str, page: int = 0) -> List[Dict]:
    chunks = []
    start = 0
    idx = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunk_text = text[start:end].strip()
        if chunk_text:
            chunks.append({
                "chunk_id": f"{doc_id}_p{page}_c{idx}",
                "doc_id": doc_id,
                "document_name": doc_id,
                "document_type": doc_type,
                "page": page,
                "text": chunk_text,
            })
            idx += 1
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks


def ingest_pdf(path: str, doc_id: str) -> List[Dict]:
    try:
        import pdfplumber
        chunks = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                chunks.extend(_chunk_text(text, doc_id, "PDF", page=i + 1))
        return chunks
    except ImportError:
        # Fallback: pypdf
        try:
            from pypdf import PdfReader
            reader = PdfReader(path)
            chunks = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                chunks.extend(_chunk_text(text, doc_id, "PDF", page=i + 1))
            return chunks
        except Exception:
            return []


def ingest_docx(path: str, doc_id: str) -> List[Dict]:
    try:
        from docx import Document
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs)
        return _chunk_text(text, doc_id, "DOCX")
    except Exception:
        return []


def ingest_pptx(path: str, doc_id: str) -> List[Dict]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        chunks = []
        for i, slide in enumerate(prs.slides):
            text = " ".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text")
            )
            chunks.extend(_chunk_text(text, doc_id, "PPTX", page=i + 1))
        return chunks
    except Exception:
        return []


def ingest_csv(path: str, doc_id: str) -> List[Dict]:
    try:
        import pandas as pd
        df = pd.read_csv(path)
        text = df.to_string(index=False)
        return _chunk_text(text, doc_id, "CSV")
    except Exception:
        return []


def ingest_xlsx(path: str, doc_id: str) -> List[Dict]:
    try:
        import pandas as pd
        chunks = []
        xl = pd.ExcelFile(path)
        for sheet in xl.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet)
            text = f"Sheet: {sheet}\n" + df.to_string(index=False)
            chunks.extend(_chunk_text(text, doc_id, "XLSX"))
        return chunks
    except Exception:
        return []


def ingest_txt(path: str, doc_id: str) -> List[Dict]:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return _chunk_text(text, doc_id, "TXT")
    except Exception:
        return []


_PARSERS = {
    ".pdf": ingest_pdf,
    ".docx": ingest_docx,
    ".pptx": ingest_pptx,
    ".csv": ingest_csv,
    ".xlsx": ingest_xlsx,
    ".xls": ingest_xlsx,
    ".txt": ingest_txt,
}


def ingest_documents(file_paths: List[str]) -> List[Dict]:
    """Ingest a list of file paths and return all chunks."""
    all_chunks = []
    for path in file_paths:
        ext = os.path.splitext(path)[1].lower()
        doc_id = os.path.basename(path)
        parser = _PARSERS.get(ext)
        if parser:
            chunks = parser(path, doc_id)
            all_chunks.extend(chunks)
    return all_chunks
